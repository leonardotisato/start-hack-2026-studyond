#!/usr/bin/env python3
"""
Studyond pitch deck — HTML → PPTX
Generates pitch-deck/studyond.pptx from the slide content.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

BASE        = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS = os.path.join(BASE, '..', 'screenshots')

W = Inches(13.33)
H = Inches(7.5)

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x09, 0x09, 0x0b)
BG2     = RGBColor(0x0f, 0x0f, 0x13)
BG3     = RGBColor(0x14, 0x14, 0x19)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
VIOLET  = RGBColor(0xA7, 0x8B, 0xFA)
BLUE    = RGBColor(0x60, 0xA5, 0xFA)
MUTED   = RGBColor(0xA1, 0xA1, 0xAA)
DIM     = RGBColor(0x52, 0x52, 0x5B)
RED     = RGBColor(0xEF, 0x44, 0x44)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
INDIGO  = RGBColor(0x63, 0x66, 0xF1)
EMERALD = RGBColor(0x22, 0xC5, 0x5E)
BORDER  = RGBColor(0x27, 0x27, 0x2a)

# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or BG

def tb(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, name='Calibri'):
    """Add a text box."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = name
    return box

def tb2(slide, lines, x, y, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, name='Calibri'):
    """Add a text box with multiple paragraphs (list of (text, color) tuples)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, (text, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = clr
        run.font.name  = name
    return box

def box(slide, x, y, w, h, fill=None, border=None, bw=Pt(1), rounded=True):
    """Add a solid shape (rounded rect or plain rect)."""
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded:
        try:
            shp.adjustments[0] = 0.04
        except Exception:
            pass
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = bw
    else:
        shp.line.fill.background()
    return shp

def pic(slide, fname, x, y, w):
    path = os.path.join(SCREENSHOTS, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w)
    else:
        # Placeholder
        b = box(slide, x, y, w, Inches(3.5), fill=BG3, border=DIM, rounded=False)
        tb(slide, f'[image: {fname}]', x, y, w, Inches(3.5),
           size=12, color=DIM, align=PP_ALIGN.CENTER)


# ── Slide 1: Hero ─────────────────────────────────────────────────────────────

def slide_hero(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)

    # Badge
    tb(sl, 'START HACK 2026',
       W/2 - Inches(1.8), Inches(1.3), Inches(3.6), Inches(0.38),
       size=10, color=DIM, align=PP_ALIGN.CENTER)

    # ThesisGPS + compass emoji  (VIOLET approximates the gradient)
    tb(sl, 'ThesisGPS  🧭',
       W/2 - Inches(4.5), H/2 - Inches(1.1), Inches(9), Inches(1.4),
       size=78, bold=True, color=VIOLET, align=PP_ALIGN.CENTER, name='Calibri')

    # powered by · Studyond
    tb(sl, 'powered by  ·  Studyond',
       W/2 - Inches(2.5), H/2 + Inches(0.5), Inches(5), Inches(0.45),
       size=17, color=MUTED, align=PP_ALIGN.CENTER)

    # Tagline
    tb(sl, 'The intelligent infrastructure for the next generation.',
       W/2 - Inches(4), H/2 + Inches(1.15), Inches(8), Inches(0.45),
       size=15, color=DIM, align=PP_ALIGN.CENTER)


# ── Slide 2: Problem ──────────────────────────────────────────────────────────

def slide_problem(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)

    # Heading (two lines, different colors)
    tb2(sl, [
        ('Every year, students start their thesis journey', WHITE),
        ('completely alone.', MUTED),
    ], Inches(1), Inches(0.75), W - Inches(2), Inches(1.4),
       size=32, bold=True, align=PP_ALIGN.CENTER)

    # 3 cards
    cards = [
        ('🧭  Decision Paralysis', 'No Structured Guidance',
         'No roadmap. No adaptive support.\nStudents are left to figure it out alone.',
         RED),
        ('🔍  Resource Scarcity', "Can't Find What They Need",
         "The right dataset or tool exists —\nbut stays invisible until it's too late.",
         AMBER),
        ('🚫  Zero Industry Impact', 'Theses Disappear',
         'A thesis ends at submission.\nNo industry connection, no career pipeline.',
         INDIGO),
    ]

    cw = Inches(3.55)
    ch = Inches(3.6)
    cy = Inches(2.4)
    gap = Inches(0.38)
    total = cw * 3 + gap * 2
    cx0 = (W - total) / 2

    for i, (tag, title, desc, accent) in enumerate(cards):
        cx = cx0 + i * (cw + gap)
        # Card bg
        box(sl, cx, cy, cw, ch, fill=BG2, border=BORDER)
        # Accent top stripe
        stripe = box(sl, cx, cy, cw, Inches(0.055), fill=accent, rounded=False)
        stripe.line.fill.background()
        # Tag
        tb(sl, tag, cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4), Inches(0.4),
           size=11, color=accent, bold=True)
        # Title
        tb(sl, title, cx + Inches(0.2), cy + Inches(0.65), cw - Inches(0.4), Inches(0.55),
           size=17, bold=True, color=WHITE)
        # Desc
        tb(sl, desc, cx + Inches(0.2), cy + Inches(1.3), cw - Inches(0.4), Inches(1.9),
           size=13, color=MUTED)


# ── Slide 3: GPS Navigator ────────────────────────────────────────────────────

def slide_gps(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)

    tb(sl, 'Thesis GPS Navigator',
       Inches(1), Inches(0.55), W - Inches(2), Inches(0.8),
       size=46, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
    tb(sl, 'A live research graph that adapts as your thesis evolves.',
       Inches(2), Inches(1.3), W - Inches(4), Inches(0.38),
       size=15, color=MUTED, align=PP_ALIGN.CENTER)

    # Browser frame
    bx = Inches(1.1); by = Inches(1.85)
    bw = W - Inches(2.2); bh = H - Inches(2.1)
    box(sl, bx, by, bw, bh, fill=BG2, border=BORDER, rounded=False)
    # URL bar
    box(sl, bx, by, bw, Inches(0.36), fill=BG3, rounded=False).line.fill.background()
    tb(sl, '● ● ●  studyond.com/thesis-gps',
       bx + Inches(0.2), by + Inches(0.04), Inches(5), Inches(0.28),
       size=9, color=DIM)
    # Screenshot
    pic(sl, 'planner.png', bx, by + Inches(0.36), bw)


# ── Slide 4: Agent System ─────────────────────────────────────────────────────

def slide_agents(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)

    tb(sl, 'The Agent System',
       Inches(1), Inches(0.55), W - Inches(2), Inches(0.8),
       size=46, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
    tb(sl, 'Claude reads your graph and proposes the next move.',
       Inches(2), Inches(1.3), W - Inches(4), Inches(0.38),
       size=15, color=MUTED, align=PP_ALIGN.CENTER)

    # Panel background
    px = Inches(0.9); py = Inches(1.85)
    pw = W - Inches(1.8); ph = H - Inches(2.15)
    box(sl, px, py, pw, ph, fill=RGBColor(0x0c, 0x0c, 0x10), border=BORDER)

    aw = Inches(2.5);  ah = Inches(1.5)
    arr_w = Inches(0.95)
    total = 3 * aw + 2 * arr_w
    ax0 = (W - total) / 2
    ay1 = Inches(2.15)

    actors1 = [
        ('👤', 'Student',     'Sets goals & reviews every proposed change',      VIOLET),
        ('🗺️', 'ThesisGPS',   'Live graph tracking research milestones',          BLUE),
        ('✦',  'Claude AI',   'Analyzes context, proposes the next path',         VIOLET),
    ]
    actor_positions = []

    for i, (emoji, name, desc, accent) in enumerate(actors1):
        ax = ax0 + i * (aw + arr_w)
        actor_positions.append(ax)
        fill_c = RGBColor(0x18, 0x14, 0x28) if accent == VIOLET else RGBColor(0x0d, 0x18, 0x2a)
        box(sl, ax, ay1, aw, ah, fill=fill_c, border=accent, bw=Pt(1.3))
        tb(sl, emoji, ax, ay1 + Inches(0.08), aw, Inches(0.48),
           size=22, align=PP_ALIGN.CENTER)
        tb(sl, name, ax, ay1 + Inches(0.52), aw, Inches(0.35),
           size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, desc, ax + Inches(0.1), ay1 + Inches(0.88), aw - Inches(0.2), Inches(0.55),
           size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

    # Arrows row 1
    arrow_labels = [('→', 'triggers'), ('↔', 'context / proposal')]
    for i, (arrow, label) in enumerate(arrow_labels):
        acx = actor_positions[i] + aw + arr_w / 2
        acy = ay1 + ah / 2
        tb(sl, arrow, acx - Inches(0.25), acy - Inches(0.22), Inches(0.5), Inches(0.38),
           size=18, color=DIM, align=PP_ALIGN.CENTER)
        tb(sl, label, acx - Inches(0.9), acy + Inches(0.18), Inches(1.8), Inches(0.26),
           size=8, color=DIM, align=PP_ALIGN.CENTER)

    # Vertical connector GPS → Scout
    vert_y = ay1 + ah
    tb(sl, '↓  when stalled',
       actor_positions[1] - Inches(0.4), vert_y + Inches(0.05), aw + Inches(0.8), Inches(0.35),
       size=9.5, color=AMBER, align=PP_ALIGN.CENTER)

    # Row 2: Scout + Network (under GPS and Claude)
    ay2 = vert_y + Inches(0.48)
    actors2 = [
        ('🔍', 'Scout Agent',      'Proactively hunts experts & resources',  AMBER),
        ('🌐', 'Studyond Network', 'Professors, companies, datasets, funding', EMERALD),
    ]
    for i, (emoji, name, desc, accent) in enumerate(actors2):
        ax = actor_positions[1 + i]
        fill_c = RGBColor(0x1c, 0x14, 0x07) if accent == AMBER else RGBColor(0x09, 0x1e, 0x13)
        box(sl, ax, ay2, aw, ah, fill=fill_c, border=accent, bw=Pt(1.3))
        tb(sl, emoji, ax, ay2 + Inches(0.08), aw, Inches(0.48),
           size=22, align=PP_ALIGN.CENTER)
        tb(sl, name, ax, ay2 + Inches(0.52), aw, Inches(0.35),
           size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, desc, ax + Inches(0.1), ay2 + Inches(0.88), aw - Inches(0.2), Inches(0.55),
           size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

    # Arrow: Scout → Network
    acx = actor_positions[1] + aw + arr_w / 2
    acy = ay2 + ah / 2
    tb(sl, '→', acx - Inches(0.25), acy - Inches(0.22), Inches(0.5), Inches(0.38),
       size=18, color=DIM, align=PP_ALIGN.CENTER)
    tb(sl, 'surfaces', acx - Inches(0.6), acy + Inches(0.18), Inches(1.2), Inches(0.26),
       size=8, color=DIM, align=PP_ALIGN.CENTER)


# ── Slide 5: Contextual Access ────────────────────────────────────────────────

def slide_contextual(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)

    tb(sl, 'Contextual Access',
       Inches(1), Inches(0.55), W - Inches(2), Inches(0.8),
       size=46, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)
    tb(sl, 'The right resource surfaces at exactly the right moment.',
       Inches(2), Inches(1.3), W - Inches(4), Inches(0.38),
       size=15, color=MUTED, align=PP_ALIGN.CENTER)

    cols = [
        ('🔬', 'Node-Triggered',       'Expert Match',
         'Instant connection to the right academic or industry expert — surfaced at the exact node where that knowledge is needed.'),
        ('🏭', 'Real Validation',       'Industry Feedback',
         '185+ companies mapped to thesis topics. Your research gets stress-tested against market needs — not just academic approval.'),
        ('🗝️', 'Always Invisible Before', 'Hidden Opportunities',
         'Datasets, funding, and networks inaccessible through traditional methods — unlocked automatically as you progress.'),
    ]

    cw = Inches(3.6);  ch = Inches(4.7)
    cy = Inches(1.85); gap = Inches(0.38)
    total = cw * 3 + gap * 2
    cx0 = (W - total) / 2

    for i, (icon, tag, keyword, desc) in enumerate(cols):
        cx = cx0 + i * (cw + gap)
        box(sl, cx, cy, cw, ch, fill=BG2, border=BORDER)
        tb(sl, icon, cx, cy + Inches(0.28), cw, Inches(0.65),
           size=38, align=PP_ALIGN.CENTER)
        tb(sl, tag, cx, cy + Inches(1.02), cw, Inches(0.3),
           size=10, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
        tb(sl, keyword, cx, cy + Inches(1.35), cw, Inches(0.6),
           size=23, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Divider
        div = box(sl, cx + cw/2 - Inches(0.22), cy + Inches(2.02),
                  Inches(0.44), Inches(0.04), fill=VIOLET, rounded=False)
        div.line.fill.background()
        tb(sl, desc, cx + Inches(0.25), cy + Inches(2.18), cw - Inches(0.5), Inches(2.2),
           size=12, color=MUTED)


# ── Slide 6: WOW ─────────────────────────────────────────────────────────────

def slide_wow(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, RGBColor(0x05, 0x05, 0x07))

    tb(sl, 'powered by  ·  thesisGPS',
       W/2 - Inches(2), Inches(0.75), Inches(4), Inches(0.4),
       size=12, color=DIM, align=PP_ALIGN.CENTER)

    # Pipeline pill
    box(sl, W/2 - Inches(4.2), Inches(1.35), Inches(8.4), Inches(0.52),
        fill=BG2, border=BORDER)
    tb(sl, 'A lonely document   →   A dynamic career pipeline',
       W/2 - Inches(4.2), Inches(1.38), Inches(8.4), Inches(0.45),
       size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # Headline
    tb(sl, 'Where research', Inches(1), Inches(2.2), W - Inches(2), Inches(1.0),
       size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, 'meets its match.', Inches(1), Inches(3.1), W - Inches(2), Inches(1.0),
       size=64, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

    tb(sl, 'Studyond — the first platform that turns every thesis into a navigated journey,\nconnected to industry, built for real impact.',
       Inches(2), Inches(4.35), W - Inches(4), Inches(0.75),
       size=14, color=MUTED, align=PP_ALIGN.CENTER)

    tb(sl, 'studyond.com',
       W/2 - Inches(1.5), Inches(5.35), Inches(3), Inches(0.42),
       size=16, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)


# ── Appendix divider ──────────────────────────────────────────────────────────

def slide_appendix(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    tb(sl, 'Appendix',
       Inches(1), H/2 - Inches(0.75), W - Inches(2), Inches(1.0),
       size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, 'Backup Slides for Q&A',
       Inches(2), H/2 + Inches(0.4), W - Inches(4), Inches(0.4),
       size=15, color=MUTED, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_hero(prs)
    slide_problem(prs)
    slide_gps(prs)
    slide_agents(prs)
    slide_contextual(prs)
    slide_wow(prs)
    slide_appendix(prs)

    out = os.path.join(BASE, 'studyond.pptx')
    prs.save(out)
    print(f'✓ Saved → {out}')

if __name__ == '__main__':
    main()
